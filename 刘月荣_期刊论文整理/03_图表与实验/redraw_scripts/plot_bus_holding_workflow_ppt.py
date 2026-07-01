from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "redrawn_figures"
OUT_DIR.mkdir(parents=True, exist_ok=True)


INK = RGBColor(36, 49, 61)
MUTED = RGBColor(99, 111, 124)
GRID = RGBColor(216, 222, 230)
BLUE = RGBColor(72, 120, 168)
BLUE_LIGHT = RGBColor(229, 238, 248)
TEAL = RGBColor(78, 141, 131)
TEAL_LIGHT = RGBColor(228, 241, 239)
AMBER = RGBColor(184, 122, 50)
AMBER_LIGHT = RGBColor(245, 234, 216)
ROSE = RGBColor(185, 87, 93)
ROSE_LIGHT = RGBColor(244, 223, 225)
GRAY = RGBColor(184, 193, 204)
GRAY_LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)


def set_line(shape, color, width=1.0, dash=None):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if dash:
        shape.line.dash_style = dash


def add_text(shape, text, font_size=8.0, bold=False, color=INK):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Songti SC"
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color


def box(slide, x, y, w, h, text, fill, edge, font_size=8.0, bold=False):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    set_line(shp, edge, 0.85)
    add_text(shp, text, font_size, bold)
    return shp


def label(slide, x, y, w, h, text, size=7.0, color=MUTED, align=PP_ALIGN.CENTER, bold=False):
    shp = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Songti SC"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return shp


def line(slide, x1, y1, x2, y2, color=INK, width=0.9, arrow=False, dash=None):
    shp = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    set_line(shp, color, width, dash)
    if arrow:
        shp.line.end_arrowhead = 3
    return shp


def elbow(slide, pts, color=INK, width=0.9, arrow=True, dash=None):
    for i, ((x1, y1), (x2, y2)) in enumerate(zip(pts, pts[1:])):
        line(slide, x1, y1, x2, y2, color, width, arrow=(arrow and i == len(pts) - 2), dash=dash)


def bus(slide, x, y, color):
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.44), Inches(0.20))
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE
    set_line(body, color, 0.8)
    for dx in (0.06, 0.18, 0.30):
        win = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + dx), Inches(y + 0.055), Inches(0.08), Inches(0.055))
        win.fill.solid()
        win.fill.fore_color.rgb = color
        win.line.fill.background()
    for dx in (0.10, 0.33):
        wheel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + dx), Inches(y + 0.16), Inches(0.06), Inches(0.06))
        wheel.fill.solid()
        wheel.fill.fore_color.rgb = INK
        wheel.line.fill.background()


def route_panel(slide):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.36), Inches(0.38), Inches(9.28), Inches(1.18))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    set_line(panel, GRID, 0.8)
    label(slide, 0.50, 0.47, 1.35, 0.18, "公交运行环境", 8.2, INK, PP_ALIGN.LEFT, True)
    y = 1.10
    line(slide, 1.18, y, 8.70, y, RGBColor(167, 176, 186), 1.1)
    for x in (1.18, 2.32, 3.46, 4.60, 5.74, 6.88, 8.02, 8.70):
        stop = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x - 0.055), Inches(y - 0.055), Inches(0.11), Inches(0.11))
        stop.fill.solid()
        stop.fill.fore_color.rgb = WHITE
        set_line(stop, RGBColor(141, 152, 165), 0.85)
    bus(slide, 1.80, 0.84, BLUE)
    bus(slide, 4.18, 0.84, TEAL)
    bus(slide, 6.40, 0.84, ROSE)
    line(slide, 2.25, 0.76, 4.10, 0.76, BLUE, 0.85, True)
    line(slide, 6.30, 0.76, 4.70, 0.76, ROSE, 0.85, True)
    label(slide, 2.90, 0.58, 0.85, 0.18, "前向时距 h⁻", 6.8, BLUE)
    label(slide, 5.18, 0.58, 0.85, 0.18, "后向时距 h⁺", 6.8, ROSE)
    label(slide, 3.34, 1.21, 0.28, 0.16, "s₃", 6.4, MUTED)
    label(slide, 5.62, 1.21, 0.28, 0.16, "s₅", 6.4, MUTED)
    label(slide, 7.90, 1.21, 0.28, 0.16, "s₇", 6.4, MUTED)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.6)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    route_panel(slide)

    box(slide, 0.60, 2.10, 1.28, 0.55, "线路参数\n客流到达率", GRAY_LIGHT, GRAY, 7.7)
    box(slide, 0.60, 2.88, 1.28, 0.55, "车辆位置\n载客状态", GRAY_LIGHT, GRAY, 7.7)
    elbow(slide, [(1.24, 2.65), (1.24, 2.88)], RGBColor(154, 164, 175), 0.8, True)

    label(slide, 2.15, 1.95, 2.0, 0.20, "强化学习驻站控制闭环", 8.4, INK, PP_ALIGN.LEFT, True)
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.12), Inches(2.10), Inches(7.04), Inches(2.76))
    frame.fill.background()
    set_line(frame, GRID, 0.85)

    y = 2.75
    items = [
        (2.36, "车辆到站\n触发事件", BLUE_LIGHT, BLUE, True),
        (3.74, "上下客服务\n计算ω", AMBER_LIGHT, AMBER, False),
        (5.12, "观测状态\nh⁻, h⁺, o/Z", BLUE_LIGHT, BLUE, False),
        (6.50, "策略网络\n输出a", ROSE_LIGHT, ROSE, True),
        (7.88, "驻站Δd\n离站运行", TEAL_LIGHT, TEAL, True),
    ]
    for x, text, fill, edge, bold in items:
        box(slide, x, y, 1.02, 0.64, text, fill, edge, 7.4, bold)
    line(slide, 1.88, 3.16, 2.36, 3.16, RGBColor(154, 164, 175), 0.85, True)
    for x in (3.38, 4.76, 6.14, 7.52):
        line(slide, x, 3.07, x + 0.36, 3.07, INK, 0.9, True)

    label(slide, 2.41, 3.48, 0.90, 0.18, "事件驱动", 6.7, BLUE)
    label(slide, 3.79, 3.48, 0.90, 0.18, "需求扰动", 6.7, AMBER)
    label(slide, 5.17, 3.48, 0.90, 0.18, "局部可观测", 6.7, BLUE)
    label(slide, 6.55, 3.48, 0.90, 0.18, "连续动作", 6.7, ROSE)
    label(slide, 7.93, 3.48, 0.90, 0.18, "环境更新", 6.7, TEAL)

    box(slide, 3.70, 4.18, 1.28, 0.52, "奖励计算\n-ω₁CV²-ω₂Tʷ", TEAL_LIGHT, TEAL, 6.8)
    box(slide, 5.35, 4.18, 1.18, 0.52, "轨迹样本\n经验缓存", GRAY_LIGHT, GRAY, 7.0)
    box(slide, 6.95, 4.18, 1.20, 0.52, "MAPPO更新\n策略参数", ROSE_LIGHT, ROSE, 7.0)
    line(slide, 4.98, 4.44, 5.35, 4.44, INK, 0.85, True)
    line(slide, 6.53, 4.44, 6.95, 4.44, INK, 0.85, True)

    # Environment/reward/update feedback uses orthogonal segments only.
    elbow(slide, [(8.39, 3.39), (8.39, 3.92), (4.34, 3.92), (4.34, 4.18)], TEAL, 0.9, True)
    elbow(slide, [(7.55, 4.18), (7.55, 3.72), (7.00, 3.72), (7.00, 3.39)], ROSE, 0.9, True)
    elbow(slide, [(8.39, 2.75), (8.39, 1.88), (8.90, 1.88), (8.90, 1.18)], TEAL, 0.9, True)
    line(slide, 8.28, 1.18, 8.90, 1.18, TEAL, 0.9, True)
    label(slide, 8.58, 2.04, 1.10, 0.22, "离站后进入下一站段", 6.8, TEAL, PP_ALIGN.LEFT)

    label(
        slide,
        2.22,
        4.92,
        6.70,
        0.20,
        "注：训练阶段利用奖励更新策略；应用阶段车辆到站后按当前策略直接给出驻站时间。",
        6.5,
        MUTED,
        PP_ALIGN.LEFT,
    )
    return prs


def main():
    prs = build_presentation()
    out = OUT_DIR / "fig_bus_holding_workflow_ppt_source.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
