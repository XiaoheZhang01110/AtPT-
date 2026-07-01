from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "redrawn_figures"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT = OUT_DIR / "fig3_improved_mappo_policy_network_wps.pptx"


W, H = 13.33, 5.0
BLUE = RGBColor(76, 115, 160)
BLUE_FILL = RGBColor(229, 237, 247)
GREEN = RGBColor(80, 135, 102)
GREEN_FILL = RGBColor(229, 241, 232)
ROSE = RGBColor(185, 89, 76)
ROSE_FILL = RGBColor(249, 232, 229)
GRAY = RGBColor(85, 96, 111)
GRAY_FILL = RGBColor(244, 246, 248)
INK = RGBColor(30, 38, 49)
LIGHT_LINE = RGBColor(200, 210, 222)


def add_textbox(slide, x, y, w, h, text, size=15, color=INK, bold=False,
                align=PP_ALIGN.CENTER, font="PingFang SC"):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
    return shape


def add_box(slide, x, y, w, h, text, fill, line, size=15, bold=False, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.35)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for idx, line_text in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "PingFang SC"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.font.bold = bold
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width=1.4):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def add_elbow(slide, x1, y1, x2, y2, xm=None, color=GRAY, width=1.35):
    if xm is None:
        xm = (x1 + x2) / 2
    add_arrow(slide, x1, y1, xm, y1, color, width)
    add_arrow(slide, xm, y1, xm, y2, color, width)
    add_arrow(slide, xm, y2, x2, y2, color, width)


def add_plus(slide, x, y, r=0.24):
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x-r), Inches(y-r), Inches(2*r), Inches(2*r))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(238, 243, 248)
    circ.line.color.rgb = BLUE
    circ.line.width = Pt(1.35)
    add_textbox(slide, x-r, y-r-0.005, 2*r, 2*r, "+", size=20, color=BLUE, bold=True)
    return circ


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

# Region headers.
add_textbox(slide, 0.36, 0.18, 2.25, 0.32, "输入信息", size=13, color=GRAY, bold=True)
add_textbox(slide, 3.25, 0.18, 3.85, 0.32, "事件图注意力增强编码", size=13, color=GRAY, bold=True)
add_textbox(slide, 8.05, 0.18, 2.85, 0.32, "策略网络输出", size=13, color=GRAY, bold=True)

# Input blocks.
add_box(slide, 0.45, 0.92, 1.90, 0.72, "当前车辆状态\ns_{i,t}", BLUE_FILL, BLUE, 15, True)
add_box(slide, 0.45, 2.23, 1.90, 0.92, "邻居事件集合\n{s_{i',t'}, a_{i',t'}, e^1, e^2}", GREEN_FILL, GREEN, 13, True)
add_textbox(slide, 0.45, 3.16, 1.90, 0.42, "∀(i',t')∈N_{i,t}, i'≠i", size=11, color=GREEN)

# Encoders.
add_box(slide, 3.05, 0.92, 1.52, 0.72, "状态编码\nFC + Tanh", BLUE_FILL, BLUE, 14)
add_box(slide, 3.05, 2.30, 1.52, 0.72, "GAT 聚合\n事件影响", GREEN_FILL, GREEN, 14)
add_box(slide, 5.35, 1.52, 1.45, 0.82, "特征融合\n[h_{i,t}, m_{i,t}]", GRAY_FILL, RGBColor(143, 161, 181), 13, True)
add_plus(slide, 4.92, 1.93)

# Actor trunk.
add_box(slide, 7.45, 1.28, 1.35, 0.72, "全连接层\nFC", ROSE_FILL, ROSE, 14)
add_box(slide, 7.45, 2.28, 1.35, 0.72, "激活函数\nTanh", ROSE_FILL, ROSE, 14)
add_box(slide, 9.45, 0.92, 1.28, 0.68, "均值\nμ", ROSE_FILL, ROSE, 14, True)
add_box(slide, 9.45, 2.30, 1.28, 0.68, "方差\nσ²", ROSE_FILL, ROSE, 14, True)
add_box(slide, 11.25, 1.52, 1.35, 0.82, "驻站动作\n输出 a_{i,t}", RGBColor(252, 239, 237), ROSE, 14, True)

# Training note, kept visually subordinate.
add_box(slide, 5.35, 3.58, 2.25, 0.58, "MAPPO 更新：优势函数 + PPO 剪切目标", RGBColor(250, 250, 250), LIGHT_LINE, 11)
add_textbox(slide, 7.70, 3.55, 2.2, 0.35, "训练阶段反馈", size=10, color=GRAY)

# Flow arrows.
add_arrow(slide, 2.35, 1.28, 3.05, 1.28, BLUE)
add_arrow(slide, 2.35, 2.66, 3.05, 2.66, GREEN)
add_elbow(slide, 4.57, 1.28, 4.70, 1.82, xm=4.75, color=BLUE)
add_elbow(slide, 4.57, 2.66, 4.70, 2.04, xm=4.75, color=GREEN)
add_arrow(slide, 5.16, 1.93, 5.35, 1.93, GRAY)
add_arrow(slide, 6.80, 1.93, 7.45, 1.64, ROSE)
add_arrow(slide, 8.12, 2.00, 8.12, 2.28, ROSE)
add_arrow(slide, 8.80, 1.64, 9.45, 1.26, ROSE)
add_arrow(slide, 8.80, 2.64, 9.45, 2.64, ROSE)
add_elbow(slide, 10.73, 1.26, 11.25, 1.88, xm=10.98, color=ROSE)
add_elbow(slide, 10.73, 2.64, 11.25, 1.98, xm=10.98, color=ROSE)
add_elbow(slide, 11.95, 2.34, 7.60, 3.87, xm=11.95, color=LIGHT_LINE, width=1.05)

# Small inline explanation labels.
add_textbox(slide, 3.00, 3.45, 1.70, 0.35, "学习邻居事件权重", size=10, color=GREEN)
add_textbox(slide, 5.00, 0.98, 2.05, 0.35, "融合局部状态与图信息", size=10, color=GRAY)
add_textbox(slide, 9.05, 3.18, 2.10, 0.35, "动作分布参数化", size=10, color=ROSE)

prs.save(OUT)
print(OUT)
